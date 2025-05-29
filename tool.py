# Lots of function aren't used anymore
# 
# # tool
import time 
import csv
import requests
from bs4 import BeautifulSoup
from urllib.parse import urlparse
from googlesearch import search
import io
import contextlib
from datetime import datetime
from ollama import Client
import asyncio
import os
from PyPDF2 import PdfReader
from pptx import Presentation
from odf.opendocument import load
from odf.text import P
import whisper
from gtts import gTTS
from pydub import AudioSegment
from pydub.playback import play
import yt_dlp
import random
import re
import json

def remove_cjk(text: str) -> str:
    # Unicode ranges for CJK characters
    cjk_ranges = [
        (0x4E00, 0x9FFF),  # CJK Unified Ideographs
        (0x3400, 0x4DBF),  # CJK Unified Ideographs Extension A
        (0x20000, 0x2A6DF),  # CJK Unified Ideographs Extension B
        (0x2A700, 0x2B73F),  # CJK Unified Ideographs Extension C
        (0x2B740, 0x2B81F),  # CJK Unified Ideographs Extension D
        (0x2B820, 0x2CEAF),  # CJK Unified Ideographs Extension E
        (0x2CEB0, 0x2EBEF),  # CJK Unified Ideographs Extension F
        (0x2F800, 0x2FA1F),  # CJK Compatibility Ideographs Supplement
        (0x3000, 0x303F),  # CJK Symbols and Punctuation
        (0x31C0, 0x31EF),  # CJK Strokes
        (0x2E80, 0x2EFF),  # CJK Radicals Supplement
    ]
    
    # Create regex pattern for CJK characters
    cjk_pattern = "[" + "".join(f"\\U{r[0]:08X}-\\U{r[1]:08X}" for r in cjk_ranges) + "]"
    cjk_regex = re.compile(cjk_pattern, re.UNICODE)
    
    # Remove CJK characters
    cleaned_text = cjk_regex.sub("", text)
    
    return cleaned_text

def get_youtube_video(query):
    search_url = f"ytsearch:{query}"  # 'ytsearch:' prefix tells yt-dlp to search
    ydl_opts = {
        'quiet': True,
        'extract_flat': True,
        'limit': 1,  # Get only the first result
    }
    
    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        info = ydl.extract_info(search_url, download=False)
    
    if 'entries' in info and len(info['entries']) > 0:
        return info['entries'][0]['url']
    
    return None

def generate_speech(text: str, lang: str = 'en', output_path: str = "output.mp3") -> str:
    # Generate speech from text using Google TTS
    tts = gTTS(text=text, lang=lang)
    tts.save(output_path)
    
    # Load the saved audio file
    audio = AudioSegment.from_file(output_path, format="mp3")
    
    # Speed up the audio by 20%
    new_speed = 1.02  # 20% faster
    audio = audio.speedup(playback_speed=new_speed)
    
    # Pitch shift (10% increase)
    semitone_increase = 1.05  # Roughly 10% pitch increase
    audio = audio._spawn(audio.raw_data, overrides={
        "frame_rate": int(audio.frame_rate * semitone_increase)
    }).set_frame_rate(audio.frame_rate)
    
    # Save the modified audio
    modified_output_path = "modified_output.mp3"
    audio.export(modified_output_path, format="mp3")
    
    return modified_output_path

CONVO_PATH = "conversation_history.json"

def save_conversation_history(history):
    with open(CONVO_PATH, "w", encoding="utf-8") as f:
        json.dump(history, f, ensure_ascii=False, indent=2)

def load_conversation_history():
    if os.path.exists(CONVO_PATH):
        with open(CONVO_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}

def transcribe_audio_files(audio_files): 
    """
    Transcribes a list of audio files using Whisper's medium model, with an optional initial prompt.
    
    :param audio_files: List of file paths to audio files.
    :param initial_prompt: Optional string to provide context or specific vocabulary for transcription.
    :return: List of transcriptions corresponding to the audio files.
    """
    model = whisper.load_model("medium")  # Load the Whisper medium model
    transcriptions = []
    
    for file in audio_files:
        result = model.transcribe(file)
        transcriptions.append(result["text"])
    
    return transcriptions

# Util
def generate_combined_prompt(system_prompt, character_prompt, knowledge_cutoff_date, character_name="assistant"):
    # Get the current date in a readable format
    current_date = datetime.now().strftime("%Y-%m-%d")

    # Replace actual newlines with the literal "\n" in the input strings
    system_prompt = system_prompt.replace("\n", "\\n")
    character_prompt = character_prompt.replace("\n", "\\n")

    # Combine the prompts and additional information
    combined_prompt = f"{system_prompt}\\nKnowledge Cutoff Date: {knowledge_cutoff_date}, Current Date: {current_date}.\\n{character_prompt}. Always stay in {character_name}."
    return combined_prompt

def analyse(file_paths, link, prompt=None):
    """
    Analyse files based on their type and path.

    :param file_paths: List of absolute file paths.
    :param link: The API link for the Ollama client.
    :return: A dictionary with results for images, documents, and audio transcriptions.
    """
    ollama_client = Client(
        host=link,
        headers={'x-some-header': 'some-value'}
    )
    print("ollama_client done")

    # Define valid extensions for images, documents, and audio
    IMAGE_EXTENSIONS = {'jpg', 'jpeg', 'png', 'bmp', 'tiff'}
    DOCUMENT_EXTENSIONS = {'txt', 'doc', 'docx', 'pdf', 'rtf', 'odt', 'pptx'}
    AUDIO_EXTENSIONS = {'mp3', 'wav', 'm4a', 'flac', 'ogg', 'aac'}

    # Separate file paths into lists based on type
    img_list = [path for path in file_paths if os.path.splitext(path)[1].strip('.').lower() in IMAGE_EXTENSIONS]
    doc_list = [path for path in file_paths if os.path.splitext(path)[1].strip('.').lower() in DOCUMENT_EXTENSIONS]
    audio_list = [path for path in file_paths if os.path.splitext(path)[1].strip('.').lower() in AUDIO_EXTENSIONS]

    print("Image Files:", img_list)
    print("Document Files:", doc_list)
    print("Audio Files:", audio_list)

    results = {}

    message_content = 'Describe the images, providing all the features from detected text to character recognition, and more.'
    if prompt:
        message_content = f"{message_content}, and provide a custom description following and in addition, found a reply to those guidelines: {prompt}"

    if img_list:
        print(img_list)
        # Process image list with Ollama
        response = ollama_client.chat(
            model='moondream', # You also need to install this model: `ollama pull moondream` While a better model could be used, I'm not focused on image quality for now
            messages=[
                {
                    'role': 'user',
                    'content': message_content,
                    'images': img_list
                }
            ],
        )
        print(response)
        results['images'] = response

    if doc_list:
        print("doc")
        document_contents = []
        for doc_path in doc_list:
            extension = os.path.splitext(doc_path)[1].strip('.').lower()
            try:
                if extension == 'pdf':
                    reader = PdfReader(doc_path)
                    pdf_text = "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
                    document_contents.append(pdf_text)
                elif extension == 'pptx':
                    presentation = Presentation(doc_path)
                    ppt_text = []
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                ppt_text.append(shape.text)
                    document_contents.append("\n".join(ppt_text))
                elif extension == 'odt':
                    odt_file = load(doc_path)
                    odt_text = []
                    for paragraph in odt_file.getElementsByType(P):
                        odt_text.append(str(paragraph))
                    document_contents.append("\n".join(odt_text))
                else:
                    with open(doc_path, 'r', encoding='utf-8') as doc_file:
                        content = doc_file.read()
                        document_contents.append(content)
            except Exception as e:
                document_contents.append(f"Error reading {doc_path}: {e}")
        results['documents'] = "\n\n".join(document_contents)
    
    if audio_list:
        print("audio")
        results['audio_transcriptions'] = transcribe_audio_files(audio_list)

    return results



def split_message(message, max_length=2000):
    """
    Splits a message into chunks of at most max_length characters.

    Args:
        message (str): The message to split.
        max_length (int): The maximum length of each chunk (default is 2000).

    Returns:
        list: A list of message chunks.
    """
    chunks = []
    while len(message) > max_length:
        split_point = message.rfind("\n", 0, max_length)
        if split_point == -1:
            split_point = max_length  # If no newline, split at max_length
        chunks.append(message[:split_point])
        message = message[split_point:].lstrip("\n")
    chunks.append(message)
    return chunks

# Tool_calls functions
# python script
def execute_script(script: str) -> str: # Could be done better ig
    """
    Executes a Python script within a restricted environment.
    
    Args:
        script (str): Python code to execute.
    
    Returns:
        str: The output or error of the script execution.
    """
    # Redirect stdout and stderr to capture output
    output = io.StringIO()
    restricted_globals = {
        "__builtins__": {
            "print": print,
            "range": range,
            "len": len,
            "str": str,
            "int": int,
            "__import__": __import__,
            # Add safe built-ins as needed
        }
    }
    restricted_locals = {}

    try:
        with contextlib.redirect_stdout(output), contextlib.redirect_stderr(output):
            exec(script, restricted_globals, restricted_locals)
    except Exception as e:
        return f"Error during execution: {e}"
    
    return output.getvalue()

# Memory manager
def save_to_csv(input_string, filename="data.csv"):
    """
    Save a string in the format "user,data" to a CSV file.

    Args:
        input_string (str): Input string in the format "user,data".
        filename (str): Name of the CSV file to save the data (default is "data.csv").

    Returns:
        None
    """
    # Parse the input string
    try:
        user, data = input_string.split(",", 1)
    except ValueError:
        raise ValueError("Input string must be in the format 'user,data'")

    # Check if the file exists
    file_exists = os.path.isfile(filename)

    # Open the file in append mode and write the data
    with open(filename, mode="a", newline="", encoding="utf-8") as csvfile:
        writer = csv.writer(csvfile)

        # Write the header if the file is new
        if not file_exists:
            writer.writerow(["User", "Data"])

        # Write the user and data
        writer.writerow([user, data])

# Depreciated, use the memory_tool function instead
def search_user_data(user, keyword='', filename="data.csv"):
    """
    Load the dataset and find all rows where the 'user' matches and 'data' contains the 'keyword'.

    Args:
        user (str): The user to search for.
        keyword (str): The keyword to look for in the data.
        filename (str): Name of the CSV file to search (default is "data.csv").

    Returns:
        list: A list of matching rows.
    """
    results = []

    # Check if the file exists
    if not os.path.isfile(filename):
        raise FileNotFoundError(f"The file '{filename}' does not exist.")

    # Open the file and search for matches
    with open(filename, mode="r", newline="", encoding="utf-8") as csvfile:
        reader = csv.DictReader(csvfile)

        for row in reader:
            if row['Data']:
                if row["User"] == user and keyword in row["Data"]:
                    results.append(row)
            else:
                if row['User'] == user:
                    results.append(row)

    return results

# Function to browse online based on query // May need to be fixed, need some test
def get_search_results(query):
    """
    Get the content of a URL if provided in the query, or the first 5 URLs (title and text content) for a given search query.

    Parameters:
    query (str): Search query string.

    Returns:
    str: A formatted string containing URL, title, and text content for each result or a message error. 
    """
    try:
        if query.startswith("http://") or query.startswith("https://"):
            # Fetch and parse content from the URL
            response = requests.get(query, timeout=10)
            if response.status_code == 200:
                soup = BeautifulSoup(response.text, 'html.parser')
                title = soup.title.string if soup.title else "No Title Found"
                text = ' '.join(soup.stripped_strings)
                return f"URL: {query}\nTitle: {title}\nContent:\n{text[:500]}...\n"
            else:
                return f"Error: Unable to fetch the URL. Status code: {response.status_code}"
        else:
            # Perform a search using DuckDuckGo HTML scraping (free method)
            search_url = f"https://html.duckduckgo.com/html/?q={query}"
            headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"}
            response = requests.get(search_url, headers=headers, timeout=10)
            
            if response.status_code != 200:
                return f"Error: Unable to perform search. Status code: {response.status_code}"

            soup = BeautifulSoup(response.text, 'html.parser')
            results = []

            for result in soup.find_all('a', class_='result__a', limit=5):
                title = result.text.strip()
                link = result['href']
                link_response = requests.get(link, timeout=10)
                link_soup = BeautifulSoup(link_response.text, 'html.parser')
                text = ' '.join(list(link_soup.stripped_strings)[:200])
                results.append(f"URL: {link}\nTitle: {title}\nContent:\n{text[:500]}...\n")

            return '\n\n'.join(results) if results else "No results found."

    except Exception as e:
        return f"Error: {str(e)}"

def get_gif_link(query: str) -> str:
    """
    Searches for a GIF on Tenor and returns a random URL from the first 5 results.
    
    :param query: The search term for the GIF.
    :return: The URL of a randomly chosen GIF from the first 5 results or an error message.
    """
    API_KEY = "AIzaSyDt8A6rjqLWoypSVbKHccQkGxUb5VjxXZw"
    url = "https://tenor.googleapis.com/v2/search"
    
    params = {
        "q": query,
        "key": API_KEY,
        "limit": 5,
        "media_filter": "minimal"
    }

    try:
        response = requests.get(url, params=params)
        response.raise_for_status()  # Raise error if response is not 200 OK
        data = response.json()

        if "results" in data and len(data["results"]) > 0:
            top_five_results = data["results"][:5]  # Get only the first 5 results
            random_gif = random.choice(top_five_results)  # Select a random GIF from these
            return random_gif["media_formats"]["gif"]["url"]

        return "No GIF found."

    except requests.exceptions.RequestException as e:
        return f"Error: {e}"